import os
import re
from docx import Document
from openai import OpenAI
import fitz  # PyMuPDF
from pptx import Presentation

# 全局客户端
client = None

# 配置API（新增超时和重试限制，防止无限转圈）
def configure_api(api_key):
    """配置API"""
    global client
    client = OpenAI(
        api_key=api_key,
        base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
        timeout=60.0,    # 每块最多等60秒，超时直接报错，绝不无限转圈
        max_retries=1    # 失败最多重试1次
    )

# 物理提取函数 - PDF文档
def extract_text_from_pdf(file_path):
    """
    使用 PyMuPDF 进行暴力物理提取，确保不丢失任何隐藏在表格或文本框中的文字。
    """
    doc = fitz.open(file_path)
    full_text = []
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        
        # 获取页面高度，执行安全裁剪（去页眉页脚 8%）
        rect = page.rect
        height = rect.height
        clip_rect = fitz.Rect(0, height * 0.08, rect.width, height * 0.92)
        
        # 💡 核心改动：使用 "text" 模式进行原始流提取
        # 这种模式不理会表格逻辑，直接按物理坐标抓取所有可见字符
        page_text = page.get_text("text", clip=clip_rect)
        
        if page_text.strip():
            full_text.append(page_text)
        else:
            # 如果裁剪后没东西，保底尝试提取全页，防止误切
            full_text.append(page.get_text("text"))
            
    doc.close()
    
    # 拼接所有页面的文本
    final_text = "\n\n".join(full_text)
    
    # 使用正则删除单独占一行的纯数字（即页码），支持多行匹配模式
    # 匹配：行首 + 可选空格 + 纯数字 + 可选空格 + 行尾和换行符
    final_text = re.sub(r'(?m)^\s*\d+\s*$\n?', '', final_text)
    
    # 把3个及以上的连续换行，统一压缩为2个换行（即保留一个标准空行）
    final_text = re.sub(r'\n{3,}', '\n\n', final_text)
    
    return final_text

# 物理提取函数 - PPTX文档
def extract_text_from_pptx(file_path):
    """
    高保真 PPTX 提取：物理遍历幻灯片中的文本框和表格
    """
    prs = Presentation(file_path)
    extracted_text = []

    for i, slide in enumerate(prs.slides):
        # 物理打标：告诉 AI 这是一张新的幻灯片
        extracted_text.append(f"\n[Slide_{i+1}]")
        
        for shape in slide.shapes:
            # 1. 提取普通文本框
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # 物理层级固化：利用 PPT 里的列表缩进级别 (level)
                        # 给文本加上缩进空格，帮助 AI 识别层级关系
                        indent = "  " * paragraph.level
                        # 如果文本以常见的序号开头，直接保留
                        extracted_text.append(f"{indent}{text}")
                        
            # 2. 提取表格（核保 PPT 里常见的费率表）
            elif shape.has_table:
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    # 将单元格内的换行符替换为空格，防止 Markdown 表格错乱
                    row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    markdown_row = "| " + " | ".join(row_data) + " |"
                    extracted_text.append(markdown_row)
                    
                    # 补齐 Markdown 表格的分割线
                    if r_idx == 0:
                        separator = "| " + " | ".join(["---"] * len(row.cells)) + " |"
                        extracted_text.append(separator)
                        
        extracted_text.append("") # 每页幻灯片结束后留空行
        
    return "\n".join(extracted_text)

# 物理提取函数 - 支持多种文件格式
def extract_text_from_file(file_path):
    """从文件中提取文本（路由分发）"""
    file_path_lower = file_path.lower()
    
    if file_path_lower.endswith('.pdf'):
        extracted_text = extract_text_from_pdf(file_path)
    elif file_path_lower.endswith('.docx'):
        extracted_text = extract_text_from_word(file_path)
    elif file_path_lower.endswith('.pptx'):
        extracted_text = extract_text_from_pptx(file_path)
    else:
        # 处理纯文本文件
        with open(file_path, 'r', encoding='utf-8') as f:
            extracted_text = f.read()
    
    # 统一后处理逻辑
    # 1. 清除单独占一行的纯数字（可能是章节装饰号或漏网的页码）
    # 匹配：行首 + 可选空格 + 纯数字 + 可选空格 + 行尾
    extracted_text = re.sub(r'(?m)^\s*\d+\s*$\n?', '', extracted_text)
    
    # 2. 压缩多余空行（将3个及以上的连续换行压缩为2个，保持排版紧凑）
    extracted_text = re.sub(r'\n{3,}', '\n\n', extracted_text)
    
    return extracted_text

# 物理提取函数 - Word文档
def extract_text_from_word(file_path):
    """
    纯物理提取：放弃底层的错误自动编号，完全依赖可见文本和正则表达式重建序号。
    专门针对排版混乱的"脏" Word 文档。
    """
    doc = Document(file_path)
    extracted_text = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            extracted_text.append("")
            continue
            
        # 核心改变：我们不再去 XML 里找 numPr，因为脏文档里的 numPr 是乱的。
        # 我们只看 text。如果 text 里真的没带序号（因为被 Word 藏在自动编号里了），
        # 我们用正则表达式检测它的特征来推断。
        
        # 探测是否是"看起来像列表，但没有物理序号"的段落
        # 通过缩进特征和简短的开头来判断
        try:
            left_indent = para.paragraph_format.left_indent
            if left_indent and left_indent.pt > 0:
                # 这是一个有缩进的段落。检查它是否已经有物理序号
                has_visible_number = re.match(r'^(?:第[一二三四五六七八九十]+[章节]|[\d一二三四五六七八九十]+[、\.]|[（\(].+[）\)]|[•●○])', text)
                
                if not has_visible_number:
                    # 如果有缩进但没有可见序号，这很可能是被隐藏的自动编号
                    # 我们强行给它加一个特定的"待补齐占位符"，而不是试图给它编号
                    text = f"[待补齐序号] {text}"
        except:
            pass # 如果无法读取段落格式，就忽略
            
        extracted_text.append(text)

    # 2. 提取表格 (保持原有逻辑)
    for table in doc.tables:
        table_text = []
        header_cells = table.rows[0].cells
        header_row = "| " + " | ".join(cell.text.strip() for cell in header_cells) + " |"
        table_text.append(header_row)
        separator = "| " + " | ".join(["---"] * len(header_cells)) + " |"
        table_text.append(separator)
        for row in table.rows[1:]:
            row_cells = row.cells
            data_row = "| " + " | ".join(cell.text.strip() for cell in row_cells) + " |"
            table_text.append(data_row)
        extracted_text.extend(table_text)
        extracted_text.append("")
        
    return "\n".join(extracted_text)

# 级别一：无损/物理清洗函数
def clean_level1(text):
    """执行无损/物理清洗"""
    # 1. 首先清理空白行和占位符
    lines = text.split('\n')
    filtered_lines = []
    
    for line in lines:
        # 跳过空白行
        if not line.strip():
            filtered_lines.append(line)
            continue
        
        # 跳过占位符
        if re.search(r'\[Image\]|\[图片\]', line):
            continue
        
        # 保留其他所有行
        filtered_lines.append(line)
    
    # 2. 清理多余的空白行
    cleaned_text = '\n'.join(filtered_lines)
    cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)
    
    return cleaned_text

# 级别二：结构化清洗函数
def clean_level2(text):
    """执行结构化清洗"""
    # 先执行级别一的清洗
    text = clean_level1(text)
    
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        # 识别标题并转化为Markdown格式
        if line.strip().startswith('#'):
            # 已经是Markdown标题，保持不变
            cleaned_lines.append(line)
        elif line.strip().isupper() and len(line.strip()) > 3:
            # 全大写的行可能是标题
            cleaned_lines.append('# ' + line.strip())
        elif line.strip().endswith(':') and len(line.strip()) < 20:
            # 以冒号结尾且长度较短的行可能是小标题
            cleaned_lines.append('## ' + line.strip())
        elif '|' in line and '---' in line:
            # 表格行，保持不变
            cleaned_lines.append(line)
        else:
            # 普通文本行
            cleaned_lines.append(line)
    
    # 清理多余的空白行
    cleaned_text = '\n'.join(cleaned_lines)
    cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)
    
    return cleaned_text

# 级别三：精粹/观点清洗函数
def clean_level3(text):
    """执行精粹/观点清洗"""
    # 先执行级别二的清洗
    text = clean_level2(text)
    
    lines = text.split('\n')
    key_points = []
    current_section = ""
    
    # 提取核心观点
    for line in lines:
        # 提取标题作为章节
        if line.strip().startswith('#'):
            current_section = line.strip()
            key_points.append(current_section)
        # 提取包含关键信息的行
        elif any(keyword in line for keyword in ['核保结论', '加费', '拒保', '标准承保', '除外责任', '延期承保']):
            key_points.append(f"- {line.strip()}")
        # 提取包含数值的行
        elif re.search(r'\d+%|\d+岁|\d+万元', line):
            key_points.append(f"- {line.strip()}")
    
    # 构建精粹内容
    cleaned_text = '\n'.join(key_points)
    
    # 添加YAML元数据
    yaml_metadata = """---
险种: 核保手册
疾病标签: [[高血压]] [[糖尿病]] [[恶性肿瘤]]
---\n\n"""
    
    return yaml_metadata + cleaned_text

# 分块函数
def chunk_text(text, chunk_size=3000, overlap=0):
    """将文本分块，按标题和逻辑结构切分，确保章节完整性"""
    chunks = []
    
    import re
    
    # 定义标题模式
    # 匹配各种标题格式：第一章、第一节、一、（一）、1.、1、（1）等
    title_regex = re.compile(r'^\s*(?:第[一二三四五六七八九十百]+[章节]|\d+(?:\.|、)|[一二三四五六七八九十百]+、|[（(][一二三四五六七八九十百]+[）)]|[（(]\d+[）)]|[①②③④⑤⑥⑦⑧⑨⑩])\s*', re.MULTILINE)
    
    # 按行分割文本
    lines = text.split('\n')
    current_chunk = []
    current_length = 0
    
    for line in lines:
        line_length = len(line)
        
        # 检查是否是标题行
        is_title = bool(title_regex.match(line))
        
        # 如果是标题行且当前块不为空，先保存当前块
        if is_title and current_chunk:
            chunk_text = '\n'.join(current_chunk)
            chunks.append(chunk_text)
            current_chunk = [line]
            current_length = line_length
        else:
            # 检查当前块大小
            if current_length + line_length > chunk_size:
                if current_chunk:
                    chunk_text = '\n'.join(current_chunk)
                    chunks.append(chunk_text)
                    current_chunk = [line]
                    current_length = line_length
            else:
                current_chunk.append(line)
                current_length += line_length + 1  # +1 for \n
    # 保存最后一个块
    if current_chunk:
        chunk_text = '\n'.join(current_chunk)
        chunks.append(chunk_text)
    
    return chunks

# AI清洗函数（新增 progress_callback 进度汇报机制）
def clean_with_ai(text, api_key, model="qwen-max", clean_level="level2", progress_callback=None):
    """使用AI清洗文本，支持进度条回调"""
    configure_api(api_key)
    chunks = chunk_text(text)
    cleaned_chunks = []
    
    # === 根据级别动态加载 System Prompt ===
    if clean_level == "level1":
        # 级别一：纯物理无损模式
        system_prompt = '''你是一台高精度文档复印机。你的任务是：
1. 原样复刻：严禁修改任何行首出现的序号（如 一、1. (1)），严禁将其转化为Markdown标题（#）。
2. 消除标记：看到 [List_LX] 或 [待补齐序号] 等提取标记时，将其移除并合理推断原序号。
3. 绝对禁止：禁止将多行列表合并为一段；禁止将中文序号改为阿拉伯数字。
4. 格式修复：修复由于底层提取导致的单词粘连和异常换行，但必须保证字数绝对无损。'''
    else:
        # 默认级别二：知识库结构化清洗模式
        system_prompt = '''你是一个专业的"知识库清洗专家"。你的任务是在【绝对不删减任何原文内容】的前提下，将纯文本排版为人类可读的格式。
清洗规则：
1. 保留原文序号：必须严格保留原文档中的 一、、1.、（1） 等所有序号和标点。
2. 禁止添加 Markdown 符号：绝对不允许在行首添加 #、-、* 等 Markdown 标记，不要画蛇添足。
3. 保持段落独立：每一个序号段落（比如 1. xxx）必须单独占一行，不要和上一段或下一段连在一起。
4. 纯净文本：输出应该是干净的人类可读文本，不要包含任何代码标记。
5. 处理占位符：如果看到 [List_L0]、[待补齐序号] 等提取标记，请将其移除并结合上下文推理正确的连续序号。'''

    for i, chunk in enumerate(chunks):
        try:
            # 💡 汇报当前进度给 Streamlit 前端
            if progress_callback:
                progress_callback(i + 1, len(chunks))
                
            context_info = ""
            if i > 0 and len(cleaned_chunks) > 0:
                previous_chunk_end = cleaned_chunks[i-1][-100:]
                context_info = f"这是文档的第 {i+1} 部分，上一部分结束于：{previous_chunk_end}\n\n请紧接上述内容继续转换，保持排版逻辑一致。\n\n"
            
            user_message = context_info + chunk
            
            # 使用新版 client 调用方式
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_message}
                ],
                temperature=0.1,
                max_tokens=4000
            )
            cleaned_chunks.append(response.choices[0].message.content)
        except Exception as e:
            # 炸出真实错误，不再静默！
            raise Exception(f"在处理第 {i+1}/{len(chunks)} 块内容时接口崩溃。错误信息：{str(e)}")
    
    # 拼接结果
    if not cleaned_chunks:
        return text
    
    final_result = cleaned_chunks[0]
    for i in range(1, len(cleaned_chunks)):
        final_result += "\n\n" + cleaned_chunks[i]
        
    import re
    if "[待补齐序号]" in final_result:
        final_result = re.sub(r'\[待补齐序号\]\s*', '', final_result)
        
    return final_result

# MD格式转换函数
def apply_md_formatting(text, level1_pattern, level2_pattern, level3_pattern, level4_pattern):
    """
    根据用户配置的正则表达式为原文添加Markdown格式
    
    Args:
        text: 干净的原文
        level1_pattern: 1级标题的正则表达式（可能包含多个模式，用 | 连接）
        level2_pattern: 2级标题的正则表达式（可能包含多个模式，用 | 连接）
        level3_pattern: 3级标题的正则表达式（可能包含多个模式，用 | 连接）
        level4_pattern: 4级标题的正则表达式（可能包含多个模式，用 | 连接）
    
    Returns:
        str: 带Markdown格式的文本
    """
    lines = text.split('\n')
    formatted_lines = []
    
    # 构建层级模式映射
    patterns = [
        (level1_pattern, '#'),
        (level2_pattern, '##'),
        (level3_pattern, '###'),
        (level4_pattern, '####')
    ]
    
    for line in lines:
        line_stripped = line.strip()
        if not line_stripped:
            formatted_lines.append(line)
            continue
        
        # 检查是否已经有 # 号
        if '#' in line:
            formatted_lines.append(line)
            continue
        
        # 检查是否匹配任何层级标识
        matched = False
        for pattern, md_format in patterns:
            if pattern:
                # 使用 re.match 匹配行首
                import re
                if re.match(pattern, line):
                    # 保持原始缩进
                    indent = re.match(r'^\s*', line).group(0)
                    # 在该行最前面插入对应的 # 和空格
                    formatted_line = f"{indent}{md_format} {line_stripped}"
                    formatted_lines.append(formatted_line)
                    matched = True
                    break
        
        if not matched:
            # 没有匹配的层级标识，保持原样
            formatted_lines.append(line)
    
    return '\n'.join(formatted_lines)

# 删除重复内容的函数
def remove_duplicates(text):
    """删除文本中的重复内容，避免删除重要信息"""
    lines = text.split('\n')
    unique_lines = []
    
    import re
    
    # 关键词列表，与verify_content函数保持一致
    keywords = ['拒保', '加费', '标准承保', '除外责任', '延期承保', '%', '岁', '万元']
    
    # 定义需要保留的行模式
    # 标题和条款开头的模式
    title_patterns = [
        r'^\s*第[一二三四五六七八九十百]+章\s*',  # 第一章、第二章等
        r'^\s*第[一二三四五六七八九十百]+节\s*',  # 第一节、第二节等
        r'^\s*[一二三四五六七八九十百]+、\s*',      # 一、二、三等
        r'^\s*[（(][一二三四五六七八九十百]+[）)]\s*',  # （一）、（二）等
        r'^\s*\d+\.\s*',                           # 1.、2.等
        r'^\s*\d+、\s*',                           # 1、2等
        r'^\s*[（(]\d+[）)]\s*',                  # （1）、（2）等
        r'^\s*[①②③④⑤⑥⑦⑧⑨⑩]\s*'                # ①、②等
    ]
    
    # 合并标题模式
    combined_title_pattern = '|'.join(title_patterns)
    
    # 记录已见过的非关键词内容
    seen_content = set()
    # 记录当前章节的标题
    current_chapter = ""
    
    for line in lines:
        line_stripped = line.strip()
        
        # 跳过空行
        if not line_stripped:
            unique_lines.append(line)
            continue
        
        # 检查是否是标题或条款开头
        is_title = bool(re.match(combined_title_pattern, line_stripped))
        
        # 检查是否包含关键词
        contains_keyword = any(keyword in line for keyword in keywords)
        
        if is_title:
            # 保留所有标题和条款开头的行
            unique_lines.append(line)
            # 更新当前章节标题
            if "章" in line_stripped or "节" in line_stripped:
                current_chapter = line_stripped
        elif contains_keyword:
            # 保留所有包含关键词的行，即使重复
            unique_lines.append(line)
        else:
            # 对于非标题、非关键词行，检查是否重复
            # 结合当前章节标题，避免不同章节的相似内容被误删
            content_key = f"{current_chapter}:{line_stripped}"
            if content_key not in seen_content:
                seen_content.add(content_key)
                unique_lines.append(line)
    
    return '\n'.join(unique_lines)

# 对账校验函数
def verify_content(original_text, cleaned_text):
    """校验内容是否一致，当关键词减少超过5%时标记为严重问题"""
    # 提取关键词
    keywords = ['拒保', '加费', '标准承保', '除外责任', '延期承保', '%', '岁', '万元']
    
    # 统计原文中的关键词数量
    original_counts = {}
    for keyword in keywords:
        original_counts[keyword] = original_text.count(keyword)
    
    # 统计清洗后文本中的关键词数量
    cleaned_counts = {}
    for keyword in keywords:
        cleaned_counts[keyword] = cleaned_text.count(keyword)
    
    # 检查是否有删减
    issues = []
    severe_issues = []
    
    for keyword, original_count in original_counts.items():
        cleaned_count = cleaned_counts.get(keyword, 0)
        if cleaned_count < original_count:
            # 计算减少百分比
            if original_count > 0:
                decrease_percent = ((original_count - cleaned_count) / original_count) * 100
                issue_msg = f"关键词 '{keyword}' 数量减少：{original_count} → {cleaned_count} ({decrease_percent:.1f}%)"
                if decrease_percent > 5:
                    severe_issues.append(issue_msg)
                else:
                    issues.append(issue_msg)
            else:
                issues.append(f"关键词 '{keyword}' 数量减少：{original_count} → {cleaned_count}")
    
    return {
        "original_counts": original_counts,
        "cleaned_counts": cleaned_counts,
        "issues": issues,
        "severe_issues": severe_issues,
        "pass": len(issues) == 0 and len(severe_issues) == 0
    }

# 差异标注函数
def annotate_differences(original_text, cleaned_text):
    """在原文中标注出清洗后可能缺失的内容"""
    # 提取关键词
    keywords = ['拒保', '加费', '标准承保', '除外责任', '延期承保', '%', '岁', '万元']
    
    # 统计原文中的关键词数量
    original_counts = {}
    for keyword in keywords:
        original_counts[keyword] = original_text.count(keyword)
    
    # 统计清洗后文本中的关键词数量
    cleaned_counts = {}
    for keyword in keywords:
        cleaned_counts[keyword] = cleaned_text.count(keyword)
    
    # 识别缺失的关键词
    missing_keywords = []
    for keyword, original_count in original_counts.items():
        cleaned_count = cleaned_counts.get(keyword, 0)
        if cleaned_count < original_count:
            missing_keywords.append(keyword)
    
    # 标注原文
    annotated_text = original_text
    
    # 在文本开头添加差异摘要
    if missing_keywords:
        summary = "⚠️ 差异标注：以下关键词在清洗过程中可能有缺失：" + ", ".join(missing_keywords) + "\n\n"
        annotated_text = summary + annotated_text
    
    # 高亮显示可能缺失的关键词
    for keyword in missing_keywords:
        # 使用Markdown高亮格式标注关键词
        annotated_text = annotated_text.replace(keyword, f"**{keyword}**")
    
    return annotated_text

# 主函数
def main():
    input_file = "test_long.txt"  # 使用长测试文件
    output_file = "output_long.md"
    
    if not os.path.exists(input_file):
        print(f"错误：文件 {input_file} 不存在")
        return
    
    print(f"开始处理文件：{input_file}")
    
    # 1. 物理提取
    raw_text = extract_text_from_file(input_file)
    print(f"提取完成，总字符数：{len(raw_text)}")
    
    # 2. 分块
    chunks = chunk_text(raw_text, chunk_size=1000, overlap=100)  # 减小分块大小以测试分块功能
    print(f"分块完成，共 {len(chunks)} 块")
    for i, chunk in enumerate(chunks):
        print(f"块 {i+1}：{len(chunk)} 字符")
    
    # 3. AI清洗
    # 使用实际的clean_with_ai函数
    api_key = "test_api_key"  # 这里需要输入实际的API key
    try:
        cleaned_text = clean_with_ai(raw_text, api_key)
        print("AI清洗完成")
    except Exception as e:
        print(f"AI清洗失败：{str(e)}")
        # 失败时使用原始文本
        cleaned_text = raw_text
    
    # 4. 使用AI清洗结果
    final_result = cleaned_text
    
    # 5. 保存结果
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(final_result)
    
    print(f"处理完成，结果已保存到：{output_file}")
    print(f"输出文件字符数：{len(final_result)}")

if __name__ == "__main__":
    main()